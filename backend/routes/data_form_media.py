from flask import Blueprint, request, jsonify
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import os
from aixplain.client import AIXplainClient
import io

data_form_media = Blueprint('data_form_media', __name__, url_prefix='/data_form_media')

# Initialize aiXplain client
AIXPLAIN_API_KEY = os.environ.get("AIXPLAIN_API_KEY")
client = AIXplainClient(api_key=AIXPLAIN_API_KEY)

# File storage directory
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'uploads')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def extract_text_from_pdf(file_path):
    text = ""
    try:
        pdf_document = fitz.open(file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text += page.get_text()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    return text

def extract_text_from_docx(file_path):
    text = ""
    try:
        document = Document(file_path)
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    return text

def analyze_text_with_aixplain(text, subject, lecture):
    """
    Analyze the extracted text using aiXplain model
    """
    try:
        # Create a prompt for the aiXplain model
        prompt = f"""
        Subject: {subject}
        Lecture: {lecture}
        
        Please analyze the following text from an educational document:
        
        {text[:2000]}  # Limiting text length for the analysis
        
        Provide a summary and key concepts covered in this document.
        """
        
        # Call aiXplain API for text analysis
        response = client.run(
            model_id="aixplain-default-text-analysis",  # Replace with your actual model ID
            data=prompt,
            task="text-generation"
        )
        
        # Extract the analysis from aiXplain's response format
        analysis = response.get("output", "Analysis could not be generated.")
        return analysis
    except Exception as e:
        print(f"Error analyzing text with aiXplain: {e}")
        return "Error during analysis"

def store_document(file_name, extracted_text, subject, date, time, lecture, analysis):
    """
    Store document metadata in a local JSON file
    """
    import json
    from datetime import datetime
    
    documents_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'data', 'documents.json')
    os.makedirs(os.path.dirname(documents_file), exist_ok=True)
    
    # Create a new document entry
    document = {
        'document_id': f"{datetime.now().timestamp()}",
        'file_name': file_name,
        'transcript': extracted_text,
        'subject': subject,
        'date': date,
        'time': time,
        'lecture': lecture,
        'analysis': analysis,
        'created_at': datetime.now().isoformat()
    }
    
    # Load existing documents or create new document list
    try:
        if os.path.exists(documents_file):
            with open(documents_file, 'r') as f:
                documents = json.load(f)
        else:
            documents = []
    except Exception:
        documents = []
    
    # Add new document and save
    documents.append(document)
    with open(documents_file, 'w') as f:
        json.dump(documents, f, indent=2)
    
    return document

@data_form_media.route('/upload_and_extract', methods=['POST'])
def upload_and_extract():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    if 'subject' not in request.form:
        return jsonify({"error": "No subject provided"}), 400
    if 'date' not in request.form:
        return jsonify({"error": "No date provided"}), 400
    if 'time' not in request.form:
        return jsonify({"error": "No time provided"}), 400
    if 'lecture' not in request.form:
        return jsonify({"error": "No lecture provided"}), 400

    file = request.files['file']
    subject = request.form['subject']
    date = request.form['date']
    time = request.form['time']
    lecture = request.form['lecture']
    
    filename = file.filename
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(file_path)

    file_extension = file.filename.split('.')[-1].lower()
    extracted_text = ""

    if file_extension == 'pdf':
        extracted_text = extract_text_from_pdf(file_path)
    elif file_extension == 'pptx':
        extracted_text = extract_text_from_pptx(file_path)
    elif file_extension == 'docx':
        extracted_text = extract_text_from_docx(file_path)
    else:
        os.remove(file_path)
        return jsonify({"error": "Unsupported file type"}), 400

    # Analyze the text using aiXplain
    analysis = analyze_text_with_aixplain(extracted_text, subject, lecture)
    
    # Store document metadata locally
    document = store_document(filename, extracted_text, subject, date, time, lecture, analysis)

    return jsonify({
        "message": "Document processed successfully",
        "document_id": document["document_id"],
        "transcript": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text,
        "analysis": analysis
    }), 200

@data_form_media.route('/get_documents', methods=['GET'])
def get_documents():
    """Get all stored documents"""
    import json
    
    documents_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'data', 'documents.json')
    
    if not os.path.exists(documents_file):
        return jsonify({"documents": []}), 200
    
    with open(documents_file, 'r') as f:
        documents = json.load(f)
    
    return jsonify({"documents": documents}), 200